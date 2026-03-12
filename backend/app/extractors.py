from __future__ import annotations

import io
from pathlib import Path

from docx import Document
from pypdf import PdfReader

try:
    from pptx import Presentation as PptxPresentation
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}


def extract_text_from_file(filename: str, data: bytes) -> str:
    suffix = Path(filename).suffix.lower()

    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(
            f"Unsupported file type '{suffix}'. Supported: PDF, DOCX, PPTX, TXT."
        )

    if suffix == ".pdf":
        reader = PdfReader(io.BytesIO(data))
        text_parts = [(page.extract_text() or "") for page in reader.pages]
        return "\n".join(text_parts).strip()

    if suffix == ".docx":
        doc = Document(io.BytesIO(data))
        text_parts = [p.text for p in doc.paragraphs]
        return "\n".join(text_parts).strip()

    if suffix == ".pptx":
        if not _PPTX_AVAILABLE:
            raise ValueError(
                "PPTX extraction requires python-pptx. "
                "Install it with: pip install python-pptx"
            )
        prs = PptxPresentation(io.BytesIO(data))
        parts: list[str] = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_parts.append(text)
            if slide_parts:
                parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_parts))
        return "\n\n".join(parts).strip()

    # .txt — try UTF-8, fall back to latin-1
    if suffix == ".txt":
        try:
            return data.decode("utf-8").strip()
        except UnicodeDecodeError:
            return data.decode("latin-1").strip()

    return ""  # unreachable given the guard above
